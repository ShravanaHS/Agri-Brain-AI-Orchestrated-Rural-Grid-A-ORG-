from pptx import Presentation

template_file = r'c:\Users\shravana HS\Desktop\amdslingshot\Project Submission Deck _ AMD Slingshot (1).pptx'
prs = Presentation(template_file)

try:
    s1 = prs.slides[0]
    s1.shapes[3].text = 'Team Details'
    s1.shapes[4].text_frame.text = 'Team name - Agri-Brain\nTeam leader name - Shravana H S\nProblem Statement - Indian Farmers Face a Triple Crisis: The \'Kilometre Gap\', Motor Burnout, and The Cloud Gap.'
except Exception as e: print(f'Error s1: {e}')

try:
    s2 = prs.slides[1]
    s2.shapes[0].text = 'Brief about the idea'
    s2.shapes[1].text = 'Agri-Brain is a Local-First Edge AI platform running entirely on an AMD Ryzen Laptop.\nIt connects to ESP32 microcontrollers in the field via MQTT, making real-time decisions without any internet. It is a brain transplant for the traditional farm.'
except Exception as e: print(f'Error s2: {e}')

try:
    s3 = prs.slides[2]
    s3.shapes[0].text = 'Opportunities & USPs'
    s3.shapes[1].text = 'Opportunities: Solving rural agricultural challenges with affordable, local AI.\n\nHow different is it from existing ideas?\nTruly Offline AI — No Cloud. Every AI model runs locally. Zero monthly subscription.\nMulti-Modal Intelligence in One Box: 5 AI engines coordinated by one gateway.\nVoice Ledger in Local Languages: spoken commands (Kannada/Hindi) logic logging.\nMotor Burnout Prevention: Acoustic AI shuts pump off before it burns.\n\nHow will it be able to solve the problem?\nBy providing a Rs 0-extra AI farm manager that works offline, speaks the local language, and protects livelihood 24/7.\n\nUSP:\n40% Water Saved. Rs 15K+ Motor repair cost avoided. 3 hrs daily walking labour saved.'
except Exception as e: print(f'Error s3: {e}')

try:
    s4 = prs.slides[3]
    s4.shapes[0].text = 'List of features offered by the solution'
    s4.shapes[1].text = '• Smart Irrigation Engine (10-grid sequential)\n• Acoustic Motor Guard (Run-dry detection)\n• Vision AI (Leaf Health, disease classification)\n• Voice + NLP Ledger (Multi-language voice parsing)\n• Soil Health Brain (Grid-level health score)\n• Gemini AI Advisor (24/7 agricultural chatbot)'
except Exception as e: print(f'Error s4: {e}')

try:
    s5 = prs.slides[4]
    s5.shapes[0].text = 'Process flow diagram or Use-case diagram'
    s5.shapes[1].text = '06:00 AM: ESP32 reads moisture, Agri-Brain auto-queues irrigation.\n07:00 AM: Vision Brain detects leaf disease. Gemini suggests treatment.\n10:00 AM: Voice command: "Irrigate Region 2". Valve opens. Ledger logs action.\n02:00 PM: Pressure drops. Acoustic Brain triggers EMERGENCY SHUTDOWN.\n06:00 PM: Soil Brain shows health scores.'
except Exception as e: print(f'Error s5: {e}')

try:
    s6 = prs.slides[6]
    s6.shapes[0].text_frame.text = 'Architecture diagram of the proposed solution'
    s6.shapes[1].text_frame.text = 'Layer 1: FIELD (ESP32, DHT22, Sensors, Solenoid Valves)\n-> MQTT\nLayer 2: COMMUNICATION (Mosquitto/HiveMQ)\n-> MQTT\nLayer 3: AMD AI GATEWAY (Irrigation, Acoustic, Vision, NLP, Soil, Gemini AI)\n->\nLayer 4: USER INTERFACE (Web Dashboard, Flutter Mobile App)'
except Exception as e: print(f'Error s6: {e}')

try:
    s7 = prs.slides[7]
    s7.shapes[0].text = 'Technologies to be used in the solution'
    s7.shapes[1].text = '• Microcontroller: ESP32 + Arduino C++\n• Simulation: Wokwi\n• Communication: MQTT / Mosquitto / HiveMQ Cloud\n• AI Gateway: Python 3\n• Acoustic AI: acoustic_brain.py\n• Vision AI: OpenCV / PIL\n• NLP/Voice: nlp_brain.py + voice_brain.py\n• Advisory AI: Google Gemini API\n• Database: SQLite (farm_ledger.db)\n• Web UI: HTML5/CSS3/Vanilla JS\n• Mobile UI: Flutter (Dart)\n• Compute: AMD Ryzen Laptop\n• PCB: KiCad'
except Exception as e: print(f'Error s7: {e}')

try:
    s8 = prs.slides[8]
    s8.shapes[0].text = 'Usage of AMD Products/Solutions'
    s8.shapes[1].text = 'All five AI models (acoustic, vision, NLP, soil, advisory) run concurrently on an AMD Ryzen Laptop CPU. Multi-core architecture enables one dedicated thread per AI brain with near-zero latency, replacing cloud AI subscriptions. Local processing speed ensures life-saving features like sub-200ms emergency motor shutdown fire on time.\nFuture scalability roadmap uses AMD ROCm for GPU-accelerated frame analysis.'
except Exception as e: print(f'Error s8: {e}')

try:
    s9 = prs.slides[9]
    s9.shapes[0].text = 'Estimated implementation cost (optional)'
    s9.shapes[1].text = 'Estimated Per-Farm Hardware Cost:\nESP32 + Solenoid Valves (x10): Rs 3,500\nMoisture + pH Sensors (x10): Rs 1,200\nLocal Wi-Fi Router: Rs 1,500\nPCB Fabrication: Rs 800\nWiring + Enclosure: Rs 500\nAMD Ryzen Laptop: Rs 0 (Existing)\n\nTotal Hardware Cost: ~Rs 7,500'
except Exception as e: print(f'Error s9: {e}')

try:
    s10 = prs.slides[10]
    s10.shapes[0].text = 'Prototype Assets (Optional)'
    s10.shapes[2].text = 'GitHub Public Repository Link:\nhttps://github.com/ShravanaHS/Agri-Brain-AI-Orchestrated-Rural-Grid-A-ORG-\n\nDemo Video Link:\n[To Be Provided]'
except Exception as e: print(f'Error s10: {e}')

out = r'c:\Users\shravana HS\Desktop\amdslingshot\docs\AgriBrain_Submission_v2.pptx'
prs.save(out)
print(f'Saved modified template to: {out}')
