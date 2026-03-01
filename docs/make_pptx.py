"""
Agri-Brain AMD Slingshot 2026 — PPTX Generator
Run: py docs/make_pptx.py
Output: docs/AgribBrain_Submission.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colours ────────────────────────────────────────────────
AMD_RED    = RGBColor(0xED, 0x1C, 0x24)
AMD_ORANGE = RGBColor(0xF7, 0x94, 0x1D)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
TEAL       = RGBColor(0x06, 0xB6, 0xD4)
PURPLE     = RGBColor(0xA8, 0x55, 0xF7)
YELLOW     = RGBColor(0xEA, 0xB3, 0x08)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG    = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG    = RGBColor(0x1E, 0x29, 0x3B)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)
VERY_DARK  = RGBColor(0x0A, 0x0A, 0x0A)

W = Inches(13.33)   # Widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helpers ────────────────────────────────────────────────
def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)

def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, bg_color=None, border_color=None, border_pt=0):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.width = Pt(border_pt)
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def label(slide, text, x, y, w=Inches(10)):
    txt(slide, text.upper(), x, y, w, Inches(0.3),
        size=9, bold=True, color=AMD_RED, align=PP_ALIGN.LEFT)

def heading(slide, text, x, y, w=Inches(11), size=36, color=WHITE):
    txt(slide, text, x, y, w, Inches(1.2),
        size=size, bold=True, color=color, align=PP_ALIGN.LEFT)

def muted_txt(slide, text, x, y, w, h, size=12):
    txt(slide, text, x, y, w, h, size=size, color=MUTED)

def card_rect(slide, x, y, w, h, accent=None):
    """Draw a dark card with optional top accent bar."""
    box(slide, x, y, w, h, bg_color=CARD_BG, border_color=RGBColor(0x33,0x41,0x55), border_pt=1)
    if accent:
        box(slide, x, y, w, Pt(4), bg_color=accent, border_pt=0)

def pill_txt(slide, text, x, y, color=None):
    if color is None:
        color = AMD_RED
    # Create a dark tinted background from the color
    r, g, b = int(color[0])//4, int(color[1])//4, int(color[2])//4
    bg_col = RGBColor(r, g, b)
    box(slide, x, y, Inches(1.4), Inches(0.28),
        bg_color=bg_col, border_color=color, border_pt=1)
    txt(slide, text, x+Inches(0.05), y, Inches(1.3), Inches(0.28),
        size=9, bold=True, color=color, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, VERY_DARK)

# Red accent stripe at top
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)

# Badge
badge = box(s, Inches(0.8), Inches(0.4), Inches(2.2), Inches(0.35), bg_color=AMD_RED)
txt(s, "AMD SLINGSHOT 2026", Inches(0.8), Inches(0.4), Inches(2.2), Inches(0.35),
    size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title
txt(s, "Agri-Brain", Inches(0.8), Inches(1.1), Inches(11), Inches(1.4),
    size=64, bold=True, color=WHITE)
txt(s, "AI-Orchestrated Rural Grid (A-ORG)", Inches(0.8), Inches(2.3), Inches(11), Inches(0.8),
    size=30, bold=True, color=AMD_ORANGE)
txt(s, 'Solving the "Kilometre Gap" with Edge AI + AMD Ryzen™',
    Inches(0.8), Inches(3.1), Inches(11), Inches(0.6),
    size=16, bold=False, color=MUTED)

# Meta cards
for i,(lbl,val) in enumerate([
    ("TEAM","Shravana H S"),("CHALLENGE","AI / Edge Computing"),
    ("TRACK","AgriTech Innovation"),("STAGE","Functional Prototype")]):
    x = Inches(0.8) + i*Inches(3.1)
    box(s, x, Inches(4.1), Inches(0.05), Inches(0.7), bg_color=AMD_RED)
    txt(s, lbl, x+Inches(0.18), Inches(4.1), Inches(2.8), Inches(0.3),
        size=8, bold=False, color=MUTED)
    txt(s, val,  x+Inches(0.18), Inches(4.4), Inches(2.8), Inches(0.4),
        size=14, bold=True, color=WHITE)

# Footer
txt(s, "AMD Slingshot 2026  |  Agri-Brain Prototype Submission  |  shravanahs97@gmail.com",
    Inches(0.8), Inches(7.0), Inches(12), Inches(0.35),
    size=8, color=RGBColor(0x33,0x41,0x55))
txt(s, "1 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 2 — Team Details
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, DARK_BG)
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "Team Details", Inches(0.8), Inches(0.3))
heading(s, "Built by One — Driven by Purpose", Inches(0.8), Inches(0.55), size=28)

# Left card — person
card_rect(s, Inches(0.8), Inches(1.4), Inches(5.8), Inches(5.5), accent=AMD_RED)
txt(s, "Shravana H S", Inches(1.0), Inches(1.6), Inches(5.4), Inches(0.6),
    size=22, bold=True, color=WHITE)
for i,(k,v) in enumerate([
    ("Role","Full-Stack AI / IoT Developer"),
    ("Email","shravanahs97@gmail.com"),
    ("GitHub","ShravanaHS"),
    ("College","Engineering Graduate"),
]):
    y = Inches(2.25) + i*Inches(0.5)
    txt(s, k+":", Inches(1.0), y, Inches(1.1), Inches(0.4), size=11, bold=True, color=AMD_ORANGE)
    txt(s, v,     Inches(2.15),y, Inches(4.2), Inches(0.4), size=11, color=WHITE)

skills = ["Python","IoT/MQTT","Edge AI","Embedded C++","Web Dev","Flutter"]
for i,sk in enumerate(skills):
    col = [AMD_RED, TEAL, GREEN, AMD_ORANGE, PURPLE, YELLOW][i]
    x = Inches(1.0) + (i%3)*Inches(1.65)
    y = Inches(4.55) + (i//3)*Inches(0.38)
    pill_txt(s, sk, x, y, col)

# Right card — stats
card_rect(s, Inches(7.0), Inches(1.4), Inches(5.5), Inches(5.5))
for i,(metric,val) in enumerate([
    ("Lines of Code","3,000+"),("AI Models Integrated","5"),("Days to Build","5")]):
    y = Inches(1.8) + i*Inches(1.6)
    txt(s, metric, Inches(7.3), y,             Inches(4.8), Inches(0.4), size=10, color=MUTED)
    txt(s, val,    Inches(7.3), y+Inches(0.35), Inches(4.8), Inches(1.0), size=40, bold=True, color=WHITE)

txt(s, "2 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 3 — Problem Statement
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, VERY_DARK)
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "Problem Statement", Inches(0.8), Inches(0.3))
heading(s, "Indian Farmers Face a Triple Crisis", Inches(0.8), Inches(0.55), size=28)

problems = [
    (AMD_RED,    "The 'Kilometre Gap'",
     "Farmers walk 5–10 km daily just to toggle irrigation valves manually,\nlosing hours of productive time every single day."),
    (AMD_ORANGE, "Motor Burnout = Financial Ruin",
     "When a pump runs dry, the motor burns. Replacing one motor costs a farmer\ntheir entire month's income — with zero warning."),
    (TEAL,       "The Cloud Gap",
     "Existing 'smart agri' solutions require 4G/internet. In rural India,\nconnectivity is not guaranteed — making cloud-dependent products useless."),
]
for i,(color,title,body) in enumerate(problems):
    y = Inches(1.5) + i*Inches(1.95)
    box(s, Inches(0.8), y, Inches(0.08), Inches(1.5), bg_color=color)
    card_rect(s, Inches(1.1), y, Inches(11.4), Inches(1.7))
    txt(s, title, Inches(1.3), y+Inches(0.15), Inches(10.8), Inches(0.5),
        size=16, bold=True, color=color)
    txt(s, body,  Inches(1.3), y+Inches(0.6),  Inches(10.8), Inches(0.9),
        size=12, color=MUTED)

txt(s, "3 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 4 — Idea / Solution Brief
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, DARK_BG)
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "About Our Idea", Inches(0.8), Inches(0.3))
heading(s, "Agri-Brain: The Farm's Autonomous Nervous System", Inches(0.8), Inches(0.55), size=24)

txt(s,
    "Agri-Brain is a Local-First Edge AI platform running entirely on an AMD Ryzen™ laptop.\n"
    "It connects to ESP32 microcontrollers in the field via MQTT, making real-time decisions\n"
    "— without any internet. It is a brain transplant for the traditional farm.",
    Inches(0.8), Inches(1.35), Inches(11.7), Inches(1.0),
    size=13, color=MUTED)

pillars = [
    (GREEN,      "🌊", "Smart Irrigation",
     "10-grid sequential system triggered by live soil moisture from ESP32 sensors."),
    (AMD_ORANGE, "🧠", "Multi-Modal AI",
     "5 AI engines running simultaneously: acoustic, vision, NLP, soil, advisory."),
    (TEAL,       "📵", "Offline-First",
     "All decisions made locally on AMD Ryzen™. Zero cloud dependency after setup."),
]
for i,(color,icon,title,body) in enumerate(pillars):
    x = Inches(0.8) + i*Inches(4.2)
    card_rect(s, x, Inches(2.6), Inches(3.9), Inches(4.4), accent=color)
    txt(s, icon,  x+Inches(0.2), Inches(2.85), Inches(3.5), Inches(0.7), size=32)
    txt(s, title, x+Inches(0.2), Inches(3.6),  Inches(3.5), Inches(0.5), size=14, bold=True, color=color)
    txt(s, body,  x+Inches(0.2), Inches(4.15), Inches(3.5), Inches(1.5), size=11, color=MUTED)

txt(s, "4 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 5 — Architecture
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, VERY_DARK)
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "Architecture Diagram", Inches(0.8), Inches(0.3))
heading(s, "System Architecture", Inches(0.8), Inches(0.55), size=28)

layers = [
    ("LAYER 1 — FIELD (ESP32 + Sensors)", GREEN,
     ["ESP32 Microcontroller","DHT22 Moisture Sensors × 10","pH Sensors × 10",
      "10× Solenoid Valves","Motor Pressure Sensor"]),
    ("LAYER 2 — COMMUNICATION (MQTT)", TEAL,
     ["HiveMQ Cloud (TLS 8883)","Local Mosquitto Broker",
      "Topics: telemetry / control / voice / ledger",
      "Bi-directional, QoS 1, Persistent"]),
    ("LAYER 3 — AMD AI GATEWAY (BRAIN)", AMD_RED,
     ["Irrigation Queue Engine","Acoustic Motor Brain",
      "Vision AI (leaf disease)","NLP + Voice Brain",
      "Soil Health Brain","Gemini AI Advisor"]),
    ("LAYER 4 — USER INTERFACE", PURPLE,
     ["Web Dashboard (HTML/CSS/JS)","Flutter Mobile App",
      "Real-time grid map","AI chat panel","Voice ledger"]),
]
for i,(title,color,items) in enumerate(layers):
    x = Inches(0.4) + i*Inches(3.25)
    card_rect(s, x, Inches(1.5), Inches(3.1), Inches(5.6), accent=color)
    txt(s, title, x+Inches(0.15), Inches(1.6), Inches(2.9), Inches(0.45),
        size=8, bold=True, color=color)
    for j,item in enumerate(items):
        y = Inches(2.2) + j*Inches(0.75)
        box(s, x+Inches(0.15), y+Inches(0.18), Inches(0.06), Inches(0.06), bg_color=color)
        txt(s, item, x+Inches(0.35), y, Inches(2.6), Inches(0.7),
            size=10, color=WHITE)

    # Arrows between layers
    if i < 3:
        ax = x + Inches(3.1) + Inches(0.07)
        txt(s, "⟷", ax, Inches(4.2), Inches(0.2), Inches(0.5),
            size=14, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.CENTER)

txt(s, "5 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 6 — Differentiation
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, DARK_BG)
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "Differentiation", Inches(0.8), Inches(0.3))
heading(s, "How Are We Different?", Inches(0.8), Inches(0.55), size=28)

diffs = [
    ("1","Truly Offline AI — No Cloud",
     "Every AI model runs locally on AMD Ryzen™. A power cut is the only thing that stops it."),
    ("2","Multi-Modal Intelligence in One Box",
     "5 specialist AI engines coordinated by one gateway — not 5 separate expensive products."),
    ("3","Voice Ledger in Local Languages",
     "Farmer says 'Added 5kg Potash' in Kannada. The system logs it — zero typing required."),
    ("4","Motor Burnout Prevention",
     "Acoustic AI listens to the pump and shuts it off before it burns. Saves ₹15,000+ per incident."),
    ("5","₹0 Extra Infrastructure",
     "Uses the farmer's existing laptop as the brain. No server rent, no monthly subscription."),
    ("6","10-Grid Sequential Irrigation",
     "Section-by-section watering maintains optimal pressure — 40% water saving vs. broadcast."),
]
for i,(num,title,body) in enumerate(diffs):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col*Inches(6.35)
    y = Inches(1.5) + row*Inches(1.85)
    card_rect(s, x, y, Inches(6.0), Inches(1.7))
    # number circle
    box(s, x+Inches(0.15), y+Inches(0.5), Inches(0.4), Inches(0.4), bg_color=AMD_RED)
    txt(s, num, x+Inches(0.15), y+Inches(0.5), Inches(0.4), Inches(0.4),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x+Inches(0.7), y+Inches(0.15), Inches(5.1), Inches(0.45),
        size=13, bold=True, color=WHITE)
    txt(s, body,  x+Inches(0.7), y+Inches(0.6),  Inches(5.1), Inches(0.9),
        size=10, color=MUTED)

txt(s, "6 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 7 — Problem → Solution Table
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, VERY_DARK)
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "Problem → Solution", Inches(0.8), Inches(0.3))
heading(s, "How Does It Solve the Problem?", Inches(0.8), Inches(0.55), size=28)

# Table header
hdr_y = Inches(1.4)
for xi,col,cw in [(Inches(0.8),"Problem",Inches(4.5)),
                   (Inches(5.5),"Agri-Brain Solution",Inches(5.8)),
                   (Inches(11.5),"Module",Inches(1.6))]:
    box(s, xi, hdr_y, cw, Inches(0.42), bg_color=AMD_RED)
    txt(s, col, xi+Inches(0.1), hdr_y+Inches(0.05), cw, Inches(0.35),
        size=10, bold=True, color=WHITE)

rows = [
    ("Farmer walks 5–10 km to toggle valves",
     "Web + Mobile app controls all 10 valves remotely from anywhere on the farm",
     "Irrigation Engine", GREEN),
    ("Pump runs dry → motor burns",
     "Acoustic AI monitors pressure PSI. Emergency shutdown triggers automatically",
     "acoustic_brain.py", AMD_ORANGE),
    ("No internet → no smart agri",
     "All AI runs on AMD Ryzen locally. MQTT via Mosquitto needs only a Wi-Fi router",
     "Local Gateway", TEAL),
    ("Crop disease spotted too late",
     "Vision AI analyses leaf images and alerts before disease spreads to entire crop",
     "vision_brain.py", PURPLE),
    ("Soil health is guesswork",
     "pH + moisture per grid → soil_brain gives actionable grid-level advice",
     "soil_brain.py", YELLOW),
    ("Literacy + language barrier",
     "Voice commands in Kannada/Hindi; NLP parses intent and logs actions hands-free",
     "voice_brain.py", AMD_RED),
]
for i,(prob,sol,mod,color) in enumerate(rows):
    row_y = Inches(1.82) + i*Inches(0.82)
    row_bg = CARD_BG if i%2==0 else RGBColor(0x16,0x21,0x32)
    box(s, Inches(0.8), row_y, Inches(12.3), Inches(0.78), bg_color=row_bg,
        border_color=RGBColor(0x33,0x41,0x55), border_pt=0.5)
    txt(s, prob, Inches(0.9), row_y+Inches(0.06), Inches(4.4), Inches(0.65), size=10, color=WHITE)
    txt(s, sol,  Inches(5.5), row_y+Inches(0.06), Inches(5.8), Inches(0.65), size=10, color=MUTED)
    txt(s, mod,  Inches(11.5),row_y+Inches(0.15), Inches(1.5), Inches(0.45), size=9, bold=True, color=color)

txt(s, "7 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 8 — USP
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, VERY_DARK)
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "Unique Selling Proposition", Inches(0.8), Inches(0.3))
heading(s, "The USP of Agri-Brain", Inches(0.8), Inches(0.55), size=28)

# Big quote box
box(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.1),
    bg_color=RGBColor(0x1A,0x0A,0x05),
    border_color=RGBColor(0x7C,0x3A,0x0D), border_pt=1)
txt(s,
    '"A ₹0-extra AI farm manager that works entirely offline,\n'
    'speaks your language, and protects your livelihood 24/7"',
    Inches(1.2), Inches(1.55), Inches(11.0), Inches(1.7),
    size=18, bold=True, color=AMD_ORANGE, align=PP_ALIGN.CENTER)

# Stats
stats = [("40%","Water Saved\nvia precision grid irrigation",GREEN),
         ("₹15K+","Motor repair cost\navoided per incident",AMD_ORANGE),
         ("3 hrs","Daily walking labour\nsaved per farmer",TEAL)]
for i,(val,desc,color) in enumerate(stats):
    x = Inches(0.8) + i*Inches(4.2)
    card_rect(s, x, Inches(3.8), Inches(3.9), Inches(2.8))
    txt(s, val,  x+Inches(0.3), Inches(4.1),  Inches(3.4), Inches(0.9),
        size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, desc, x+Inches(0.3), Inches(5.0), Inches(3.4), Inches(1.0),
        size=11, color=MUTED, align=PP_ALIGN.CENTER)

txt(s, "8 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 9 — Feature List
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, DARK_BG)
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "Feature List", Inches(0.8), Inches(0.3))
heading(s, "Complete Feature Set", Inches(0.8), Inches(0.55), size=28)

features = [
    (GREEN,      "Smart Irrigation Engine",
     "• 10-grid sequential queue with configurable duration\n• Auto-trigger when humidity < 50%\n• Manual override via web/mobile/voice\n• 2-hour ESP32 hard safety timeout"),
    (AMD_ORANGE, "Acoustic Motor Guard",
     "• Real-time pump pressure PSI monitoring\n• Run-dry detection via pressure anomaly\n• Emergency shutdown — all valves OFF\n• Instant alert to dashboard + mobile"),
    (PURPLE,     "Vision AI — Leaf Health",
     "• Early/Late Blight detection in tomatoes\n• Leaf Miner, Mite, Aphid pest detection\n• Nutrient deficiency via leaf colour\n• Processed locally on AMD Ryzen™"),
    (YELLOW,     "Voice + NLP Ledger",
     "• Voice input in Kannada, Hindi, English\n• NLP parses action, material, quantity, grid\n• Persistent SQLite ledger (farm_ledger.db)\n• Hands-free irrigation via voice command"),
    (TEAL,       "Soil Health Brain",
     "• Per-grid pH + moisture analysis\n• Actionable advice (Add lime / Reduce water)\n• Grid-level health score mapping"),
    (AMD_RED,    "Gemini AI Advisor",
     "• 24/7 agricultural advisory chatbot\n• Crop rotation, water & pest advice\n• Integrated into web dashboard AI chat"),
]
for i,(color,title,body) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col*Inches(6.35)
    y = Inches(1.45) + row*Inches(1.95)
    card_rect(s, x, y, Inches(6.0), Inches(1.8), accent=color)
    txt(s, title, x+Inches(0.2), y+Inches(0.2), Inches(5.6), Inches(0.45),
        size=13, bold=True, color=color)
    txt(s, body,  x+Inches(0.2), y+Inches(0.65), Inches(5.6), Inches(1.1),
        size=10, color=MUTED)

txt(s, "9 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 10 — Use-Case / Process Flow
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, VERY_DARK)
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "Use-Case & Process Flow", Inches(0.8), Inches(0.3))
heading(s, "A Day in the Life of Agri-Brain", Inches(0.8), Inches(0.55), size=28)

# vertical timeline line
box(s, Inches(1.55), Inches(1.4), Inches(0.05), Inches(5.8),
    bg_color=RGBColor(0x33,0x41,0x55))

steps = [
    (AMD_RED,    "06:00 AM",
     "Morning Sensor Report — ESP32 reads moisture + pH. Agri-Brain detects Grid 3 at 32% humidity → auto-queues irrigation."),
    (AMD_ORANGE, "07:00 AM",
     "Farmer uploads leaf photo → Vision Brain detects Early Blight Stage 1. Gemini suggests treatment protocol."),
    (GREEN,      "10:00 AM",
     "Voice command in Kannada: 'Irrigate Region 2 for 15 min'. NLP Brain parses it. Valve opens. Ledger logs action."),
    (PURPLE,     "02:00 PM",
     "Pressure PSI drops. Acoustic Brain triggers EMERGENCY SHUTDOWN. All 10 valves close. Motor is saved."),
    (TEAL,       "06:00 PM",
     "Daily report: Soil Brain shows per-grid health scores. Dashboard summarises water used, alerts raised, entries logged."),
]
for i,(color,time,body) in enumerate(steps):
    y = Inches(1.4) + i*Inches(1.15)
    # dot
    box(s, Inches(1.42), y+Inches(0.1), Inches(0.28), Inches(0.28), bg_color=color)
    txt(s, time, Inches(1.85), y,              Inches(1.5), Inches(0.35),
        size=10, bold=True, color=color)
    txt(s, body, Inches(1.85), y+Inches(0.35), Inches(10.9), Inches(0.75),
        size=10, color=MUTED)

txt(s, "10 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 11 — Technology Stack
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, DARK_BG)
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "Technologies Used", Inches(0.8), Inches(0.3))
heading(s, "Technology Stack", Inches(0.8), Inches(0.55), size=28)

# Header
hdr_y = Inches(1.35)
for xi,col,cw in [(Inches(0.8),"Layer",Inches(2.3)),
                   (Inches(3.3),"Technology",Inches(3.8)),
                   (Inches(7.3),"Purpose",Inches(5.7))]:
    box(s, xi, hdr_y, cw, Inches(0.38), bg_color=AMD_RED)
    txt(s, col, xi+Inches(0.1), hdr_y+Inches(0.04), cw, Inches(0.32),
        size=9, bold=True, color=WHITE)

tech_rows = [
    ("Microcontroller","ESP32 + Arduino C++","Field sensor reading, valve control, DHT22 sensors"),
    ("Simulation","Wokwi","Zero-hardware ESP32 prototype validation"),
    ("Communication","MQTT / HiveMQ Cloud (TLS 8883)","Bi-directional real-time telemetry & control"),
    ("AI Gateway","Python 3 / paho-mqtt","Central orchestrator running all AI brains"),
    ("Acoustic AI","acoustic_brain.py","Motor pressure anomaly detection & emergency shutdown"),
    ("Vision AI","Python + OpenCV / PIL","Leaf disease and pest classification from images"),
    ("NLP / Voice","nlp_brain.py + voice_brain.py","Multi-language voice command parsing"),
    ("Advisory AI","Google Gemini API","Contextual 24/7 agricultural advice chatbot"),
    ("Database","SQLite (farm_ledger.db)","Offline-first farm activity ledger"),
    ("Web UI","HTML5 / CSS3 / Vanilla JS","Dashboard: irrigation control, grid map, AI chat"),
    ("Mobile UI","Flutter (Dart)","Cross-platform mobile control app"),
    ("Compute","AMD Ryzen™ Laptop","All AI inference — replaces cloud entirely"),
    ("PCB","KiCad","Custom PCB design for field deployment"),
]
for i,(layer,tech,purpose) in enumerate(tech_rows):
    row_y = Inches(1.73) + i*Inches(0.405)
    row_bg = CARD_BG if i%2==0 else VERY_DARK
    box(s, Inches(0.8), row_y, Inches(12.1), Inches(0.38), bg_color=row_bg)
    txt(s, layer,   Inches(0.9), row_y,          Inches(2.2), Inches(0.36), size=9, bold=True,  color=AMD_ORANGE)
    txt(s, tech,    Inches(3.3), row_y,          Inches(3.8), Inches(0.36), size=9, bold=False, color=WHITE)
    txt(s, purpose, Inches(7.3), row_y,          Inches(5.5), Inches(0.36), size=9, bold=False, color=MUTED)

txt(s, "11 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 12 — AMD Usage
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, RGBColor(0x14,0x05,0x05))
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "AMD Products / Solutions", Inches(0.8), Inches(0.3))
heading(s, "Why AMD Powers This?", Inches(0.8), Inches(0.55), size=28)

amd_points = [
    ("AMD Ryzen™ — The Edge AI Brain",
     "All five AI models (acoustic, vision, NLP, soil, advisory) run concurrently on AMD Ryzen™ CPU.\n"
     "Multi-core architecture enables one dedicated thread per AI brain with near-zero latency.\n"
     "This replaces an ₹8,000/month cloud AI subscription with a one-time local setup."),
    ("Local Processing → Zero Latency Safety Responses",
     "When Acoustic Brain detects motor anomaly, EMERGENCY SHUTDOWN fires in <200 ms.\n"
     "A cloud-based system takes 2–5 seconds (round-trip). Those seconds mean a burnt motor.\n"
     "AMD's local processing speed makes this life-saving feature viable."),
    ("Scalability Path: AMD ROCm™ (Future)",
     "Vision AI is designed for GPU acceleration. Future versions will use AMD Radeon™ + ROCm™\n"
     "to run ResNet / EfficientNet crop-disease models with sub-100 ms inference.\n"
     "This enables real-time field video analysis at production scale."),
]
for i,(title,body) in enumerate(amd_points):
    y = Inches(1.4) + i*Inches(1.9)
    box(s, Inches(0.8), y, Inches(11.7), Inches(1.75),
        bg_color=RGBColor(0x1A,0x05,0x05),
        border_color=RGBColor(0x7C,0x20,0x20), border_pt=1)
    # left accent bar
    box(s, Inches(0.8), y, Inches(0.08), Inches(1.75), bg_color=AMD_RED)
    txt(s, title, Inches(1.05), y+Inches(0.12), Inches(11.2), Inches(0.45),
        size=14, bold=True, color=AMD_ORANGE)
    txt(s, body,  Inches(1.05), y+Inches(0.55), Inches(11.2), Inches(1.1),
        size=10, color=MUTED)

txt(s, "12 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# SLIDE 13 — Cost + Roadmap + Links
# ──────────────────────────────────────────────────────────
s = blank_slide(prs)
bg(s, DARK_BG)
box(s, 0, 0, W, Inches(0.06), bg_color=AMD_RED)
label(s, "Cost & Roadmap", Inches(0.8), Inches(0.3))
heading(s, "Implementation Cost & Next Steps", Inches(0.8), Inches(0.55), size=24)

# Cost Table (left)
txt(s, "Estimated Per-Farm Hardware Cost", Inches(0.8), Inches(1.35), Inches(5.8), Inches(0.35),
    size=11, bold=True, color=MUTED)
cost_items = [
    ("ESP32 + Solenoid Valves (×10)", "₹3,500"),
    ("Moisture + pH Sensors (×10)",   "₹1,200"),
    ("Local Wi-Fi Router",            "₹1,500"),
    ("PCB Fabrication (Custom)",      "₹800"),
    ("Wiring + Enclosure",            "₹500"),
    ("AMD Ryzen™ Laptop (co-op)",     "₹0"),
]
for i,(item,cost) in enumerate(cost_items):
    y = Inches(1.75) + i*Inches(0.58)
    row_bg = CARD_BG if i%2==0 else VERY_DARK
    box(s, Inches(0.8), y, Inches(5.8), Inches(0.54), bg_color=row_bg)
    txt(s, item, Inches(0.95), y+Inches(0.08), Inches(3.9), Inches(0.38), size=10, color=WHITE)
    txt(s, cost, Inches(5.2),  y+Inches(0.08), Inches(1.2), Inches(0.38), size=11, bold=True,
        color=GREEN, align=PP_ALIGN.RIGHT)
# Total
box(s, Inches(0.8), Inches(5.27), Inches(5.8), Inches(0.6), bg_color=AMD_RED)
txt(s, "Total Hardware Cost", Inches(0.95), Inches(5.35), Inches(3.0), Inches(0.42),
    size=12, bold=True, color=WHITE)
txt(s, "~₹7,500", Inches(4.5), Inches(5.35), Inches(1.9), Inches(0.42),
    size=16, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# Roadmap (right)
txt(s, "Prototype → Production Roadmap", Inches(7.2), Inches(1.35), Inches(5.8), Inches(0.35),
    size=11, bold=True, color=MUTED)
roadmap = [
    (True,  "Phase 1: Wokwi + Cloud Broker",    "10-grid simulation, AI gateway, HiveMQ integration"),
    (True,  "Phase 2: Premium Web Dashboard",    "Grid map, real-time telemetry, AI chat, voice ledger"),
    (True,  "Phase 3: Custom PCB Design",        "KiCad PCB for field-ready hardware deployment"),
    (False, "Phase 4: Physical Field Pilot",     "Deploy on real farm, integrate sensors, iterate AI models"),
    (False, "Phase 5: AMD ROCm™ GPU Inference",  "Real-time video crop analysis with AMD Radeon™"),
]
for i,(done,title,desc) in enumerate(roadmap):
    y = Inches(1.75) + i*Inches(1.05)
    dot_color = GREEN if done else RGBColor(0x33,0x41,0x55)
    box(s, Inches(7.2), y+Inches(0.08), Inches(0.18), Inches(0.18), bg_color=dot_color)
    status = "✓ Done" if done else "Upcoming"
    status_color = GREEN if done else MUTED
    txt(s, title,  Inches(7.55), y,              Inches(4.5), Inches(0.38), size=12, bold=True, color=WHITE)
    txt(s, status, Inches(11.0), y+Inches(0.04), Inches(1.5), Inches(0.32), size=9,  bold=True, color=status_color)
    txt(s, desc,   Inches(7.55), y+Inches(0.38), Inches(5.1), Inches(0.55), size=9,  color=MUTED)
    # Connector line
    if i < 4:
        box(s, Inches(7.29), y+Inches(0.26), Inches(0.02), Inches(0.8),
            bg_color=RGBColor(0x33,0x41,0x55))

# Links at bottom
box(s, Inches(0.8), Inches(6.2), Inches(5.8), Inches(0.8), bg_color=AMD_RED)
txt(s, "🔗  github.com/ShravanaHS/Agri-Brain-AI-Orchestrated-Rural-Grid-A-ORG-",
    Inches(1.0), Inches(6.37), Inches(10.0), Inches(0.45),
    size=11, bold=True, color=WHITE)

txt(s, "13 / 13", Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x33,0x41,0x55), align=PP_ALIGN.RIGHT)

# ──────────────────────────────────────────────────────────
# Save
# ──────────────────────────────────────────────────────────
out = r"c:\Users\shravana HS\Desktop\amdslingshot\docs\AgriBrain_Submission.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
